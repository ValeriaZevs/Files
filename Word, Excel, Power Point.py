Python 3.9.0 (tags/v3.9.0:9cf6752, Oct  5 2020, 15:34:40) [MSC v.1927 64 bit (AMD64)] on win32
Type "help", "copyright", "credits" or "license()" for more information.
>>> from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import sys

document = Document()

data = list(map(lambda x: x.replace('\n', ''), sys.stdin.readlines()))

paragraph = document.add_paragraph()
run = paragraph.add_run(data[0])
run.italic = True
font = run.font
font.name = 'Arial'
font.size = Pt(11)
par = paragraph.paragraph_format
par.alignment = WD_ALIGN_PARAGRAPH.RIGHT

for i in data[1:-2]:
    paragraph = document.add_paragraph()
    run = paragraph.add_run(i)
    font = run.font
    font.name = 'Times New Roman'
    font.size = Pt(12)
    paragraph_format = paragraph.paragraph_format
    paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    paragraph_format.first_line_indent = Inches(0.5)


paragraph = document.add_paragraph()
run = paragraph.add_run(data[-2])
paragraph_format = paragraph.paragraph_format
paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT

paragraph = document.add_paragraph()
run2 = paragraph.add_run(data[-1])
run2.bold = True
paragraph_format = paragraph.paragraph_format
paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT

document.save('letter.docx')
-----------------------------------------
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
import sys

document = Document()

data = list(map(lambda x: x.replace('\n', ''), sys.stdin.readlines()))

head = document.add_heading('Blood test', 0)
paragraph_format = head.paragraph_format
paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

table = document.add_table(len(data) + 1, 3)
heading = ("indicator", "norm", "value")
head_row = table.rows[0].cells
for i in range(len(heading)):
    p = head_row[i].paragraphs[0]
    p.add_run(heading[i]).bold = True
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
for i in range(len(data)):
    cell = table.cell(i + 1, 0)
    cell.text = str(data[i])
table.style = 'Table Grid'
document.save('analysis.docx')
-------------------------------------------
from docx import Document
import sys

document = Document()

data = list(map(lambda x: x.replace('\n', ''), sys.stdin.readlines()))

head = document.add_heading(data[0], 0)
paragraph_format = head.paragraph_format

for i in data:
    if i[0] != ' ':
        par = document.add_paragraph(i).style = 'List Bullet'
    elif i[0] == ' ' and i[1] == ' ' and i[2] != ' ':
        document.add_paragraph(i, style='List Bullet 2')
    else:
        document.add_paragraph(i, style='List Bullet 3')
document.save('systematization.docx')
--------------------------------------------
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches
from docxtpl import DocxTemplate


def template(name):
    document = Document()

    p = document.add_heading(f"Самостоятельная работа №" + '{{ work }}', 1)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2 = document.add_heading(f"Вариант №" + '{{ variant }}', 2)
    p2.alignment = WD_ALIGN_PARAGRAPH.LEFT
    p3 = document.add_paragraph(f"Задание 1\n" + '{{ task1 }}')
    p3.style = "List Number"
    p3 = document.add_paragraph(f"Задание 2\n" + '{{ task2 }}')
    p3.style = "List Number"
    p3 = document.add_paragraph(f"Задание 3\n" + '{{ task3 }}')
    p3.style = "List Number"
    p6 = document.add_paragraph()
    style = p6.add_run(f"Выполнил: " + '{{ name }}')
    p6f = p6.paragraph_format
    p6f.first_line_indent = Inches(0.5)
    style.italic = True

    document.save(name)


def prepare_doc(nameOfTemolate, **kwargs):
    doc = DocxTemplate(nameOfTemolate)
    doc.render(kwargs)
    doc.save("res.docx")
-------------------------------------------
import xlsxwriter
import sys

data = list(map(lambda x: x.replace('\t', ' '), sys.stdin.readlines()))
data = [i.replace('\n', '') for i in data]
p = int(data[-1])
datan = []
for i in data[:-1]:
    j = i.split(' ')
    r = (j[0] + ' ' + j[1], ' '.join(j[2:-1]), int(j[-1]))
    datan.append(r)
data = datan

workbook = xlsxwriter.Workbook('statement.xlsx')
worksheet = workbook.add_worksheet('Sheet_1')
cell_format = workbook.add_format()
cell_format.set_bold()

headings = 'employee, department, work, payment'.split(', ')
for row, item in enumerate(headings):
    worksheet.write(0, row, headings[row], cell_format)

for row, (man, dep, pay) in enumerate(data):
    worksheet.write(row + 1, 0, man)
    worksheet.write(row + 1, 1, dep)
    worksheet.write(row + 1, 2, str(pay))
    worksheet.write(row + 1, 3, pay * p)


worksheet = workbook.add_worksheet('Sheet_2')
cell_format = workbook.add_format()
cell_format.set_bold()

headings = 'department, total work, total cost'.split(', ')
for row, item in enumerate(headings):
    worksheet.write(0, row, headings[row], cell_format)

work = sorted(list(set([i[1] for i in data])))
for row, item in enumerate(work):
    worksheet.write(row + 1, 0, item)
points = [int(i[-1]) for i in data]
for row, i in enumerate(work):
    r = 0
    for j in data:
        if j[1] == i:
            r += j[-1]
    worksheet.write(row + 1, 1, r)
    worksheet.write(row + 1, 2, r * p)

workbook.close()
-----------------------------------------------
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Коллекции в Python"
subtitle.text = "Самое важное"

empty = prs.slides.add_slide(prs.slide_layouts[3])
title = empty.shapes.title
sub = empty.placeholders[1]
sub2 = empty.placeholders[2]
title.text = "Виды коллекций"

tf = sub.text_frame
p2 = tf.add_paragraph()
p2.text = 'Неизменяемые коллекции'
p2.level = 0
p2.font.bold = True
p3 = tf.add_paragraph()
p3.text = 'Строки'
p3.level = 1
p4 = tf.add_paragraph()
p4.text = 'Кортежи'
p4.level = 1

tf2 = sub2.text_frame
p2 = tf2.add_paragraph()
p2.text = "Изменяемые коллекции"
p2.level = 0
p2.font.bold = True
p2_2 = tf2.add_paragraph()
p2_2.text = "Списки"
p2_2.level = 1
p2_3 = tf2.add_paragraph()
p2_3.text = "Множества"
p2_3.level = 1
p2_4 = tf2.add_paragraph()
p2_4.text = "Словари"
p2_4.level = 1

left = Inches(1)
top = Inches(5.5)
txBox = empty.shapes.add_textbox(left, top, 0, 0)
tf3 = txBox.text_frame
p3 = tf3.add_paragraph()
run = p3.add_run()
run.text = "Документация"
run.font.name = "Courier New"
run.font.color.rgb = RGBColor(0xff, 0x7f, 0x50)
run.hyperlink.address = "https://python-pptx.readthedocs.io/_/downloads/en/latest/pdf/"

empty = prs.slides.add_slide(prs.slide_layouts[6])
left = Inches(2.0)
top = Inches(6.0)
width = Inches(8)
height = Inches(1)

txt_box = empty.shapes.add_textbox(left, top, width, height)
txt_frame = txt_box.text_frame
p = txt_frame.add_paragraph()
left = Inches(1.0)
top = Inches(0.5)
width = Inches(8)
height = Inches(1)
txt_box2 = empty.shapes.add_textbox(left, top, width, height)
txt_frame2 = txt_box2.text_frame
p2 = txt_frame2.add_paragraph()
run = p2.add_run()
run.text = "Свойства коллекций"
run.font.bold = True
run.font.size = Pt(40)
p2.alignment = PP_ALIGN.CENTER

run = p.add_run()
run.text = "Все коллекции поддерживают общие действия:\n"
run2 = p.add_run()
run2.text = "перебор элементов в цикле,\n"
run2.font.color.rgb = RGBColor(0, 127, 255)
run2.font.name = "Arial"
run2.font.size = Pt(18)
run2.font.italic = True
run = p.add_run()
run.text = "получение одного элемента."

left = Inches(1.2)
top = Inches(1.9)
width = Inches(8)
pic = empty.shapes.add_picture("collections.png", left, top, width)

prs.save('test.pptx')